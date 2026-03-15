"""
PPT(X) 파싱 서비스 – python-pptx 기반
- 텍스트 추출: 텍스트박스, 표, 도형, 슬라이드 노트, 숨겨진 슬라이드, 하이퍼링크
- 슬라이드 이미지 렌더링: 내장 이미지 추출 (OCR / Vision AI 입력용)
- 메타데이터 추출
"""
from __future__ import annotations
import io
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from PIL import Image

from core.config import get_logger

logger = get_logger("ppt_service")


def _iter_group_shapes(group_shape):
    """그룹 도형을 재귀적으로 순회해 모든 하위 shape 반환 (중첩 그룹 지원)"""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        for child in group_shape.shapes:
            if child.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _iter_group_shapes(child)  # 재귀
            else:
                yield child
    except Exception:
        pass


@dataclass
class SlideData:
    slide_number: int          # 1-based
    text:         str  = ""    # 모든 텍스트 합산
    text_items:   list = field(default_factory=list)  # [(source, text), ...]
    images:       list = field(default_factory=list)  # image dicts (crop_visible 포함)
    hyperlinks:   list = field(default_factory=list)  # [url, ...]
    is_hidden:    bool = False


class PPTService:
    def __init__(self, path: Path):
        self.path         = path
        self._prs         = None
        self.total_slides = 0
        self.metadata:    dict = {}
        self.is_pptx:     bool = True

    # ── 열기 ─────────────────────────────────────────────────
    def open(self) -> bool:
        try:
            from pptx import Presentation
            self._prs         = Presentation(str(self.path))
            self.total_slides = len(self._prs.slides)
            # slides[idx] 인덱스 접근이 일부 파일에서 AttributeError 발생하는 버그 대비
            # → 슬라이드를 미리 list로 캐싱해 안전하게 접근
            self._slides_list = list(self._prs.slides)
            self.metadata     = self._extract_metadata()
            # presentation.xml 의 firstSlideNum 읽기
            # PowerPoint [디자인 → 슬라이드 크기 → 시작 슬라이드 번호] 에서 사용자가 설정한 값
            # 없으면 기본값 1 (MS PowerPoint 표준)
            self.first_slide_num = self._read_first_slide_num()
            logger.info(f"PPTX 열기: {self.total_slides}슬라이드, firstSlideNum={self.first_slide_num}")
            return True
        except Exception as e:
            logger.error(f"PPTX 열기 실패: {e}")
            return False

    def _read_first_slide_num(self) -> int:
        """presentation.xml 에서 firstSlideNum 속성을 읽어 반환한다.
        없으면 1 (PowerPoint 기본값) 반환.
        """
        try:
            import zipfile, re
            with zipfile.ZipFile(str(self.path), 'r') as z:
                xml = z.read('ppt/presentation.xml').decode('utf-8', errors='replace')
            m = re.search(r'firstSlideNum="(\d+)"', xml)
            return int(m.group(1)) if m else 1
        except Exception:
            return 1

    def close(self):
        self._prs = None

    def __enter__(self):  self.open();  return self
    def __exit__(self, *_): self.close()

    # ── 메타데이터 ────────────────────────────────────────────
    def _extract_metadata(self) -> dict:
        out = {}
        try:
            cp = self._prs.core_properties
            fields = [
                ("author",           cp.author),
                ("last_modified_by", cp.last_modified_by),
                ("company",          getattr(cp, "company", "")),
                ("title",            cp.title),
                ("subject",          cp.subject),
                ("keywords",         cp.keywords),
            ]
            for k, v in fields:
                if v and str(v).strip():
                    out[k] = str(v).strip()
        except Exception as e:
            logger.warning(f"메타데이터 추출 실패: {e}")
        return out

    # ── 슬라이드 추출 ─────────────────────────────────────────
    def extract_slide(self, idx: int) -> SlideData:
        """0-based index"""
        # slides[idx] 인덱스 접근 버그 대비 → 캐싱된 list 사용
        slide      = self._slides_list[idx]
        slide_num  = idx + 1
        text_items = []
        images     = []
        hyperlinks = []

        # 숨겨진 슬라이드 여부
        is_hidden = False
        try:
            from pptx.oxml.ns import qn
            show = slide._element.get("show")
            is_hidden = (show == "0")
        except Exception:
            pass

        # ── 모든 shape 순회 ──────────────────────────────────
        for shape in slide.shapes:
            # 텍스트박스 / 도형
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    # ★ para.text 우선 사용: runs가 없어도 전체 텍스트 반환
                    #   (runs만 순회하면 runs=[] 단락의 텍스트가 완전 누락됨)
                    line = para.text or ""

                    # 하이퍼링크는 runs에서만 추출 가능
                    for run in para.runs:
                        try:
                            if run.hyperlink and run.hyperlink.address:
                                hyperlinks.append(run.hyperlink.address)
                        except Exception:
                            pass

                    if line.strip():
                        src = "note" if shape.shape_type == 13 else "textbox"
                        text_items.append((src, line.strip()))

            # 표(Table)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        ct = cell.text.strip()
                        if ct:
                            text_items.append(("table", ct))
                            # 셀 하이퍼링크
                            try:
                                for para in cell.text_frame.paragraphs:
                                    for run in para.runs:
                                        if run.hyperlink and run.hyperlink.address:
                                            hyperlinks.append(run.hyperlink.address)
                            except Exception:
                                pass

            # 이미지 추출 (Picture shape)
            try:
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_blob = shape.image.blob
                    img_ext  = shape.image.ext          # "jpeg", "png", …
                    pil_img  = Image.open(io.BytesIO(img_blob))
                    w, h     = pil_img.size

                    # 크롭 정보 추출: python-pptx의 crop 속성 (0.0 ~ 1.0 비율)
                    # crop_left/right/top/bottom: 잘려나간 비율 (0이면 크롭 없음)
                    crop_info = {"left": 0.0, "top": 0.0, "right": 0.0, "bottom": 0.0}
                    is_cropped = False
                    try:
                        cl = getattr(shape, 'crop_left',   0) or 0
                        ct = getattr(shape, 'crop_top',    0) or 0
                        cr = getattr(shape, 'crop_right',  0) or 0
                        cb = getattr(shape, 'crop_bottom', 0) or 0
                        # python-pptx crop 값은 914400분의 1 단위 (EMU ratio)이므로 나누기
                        # 실제로는 비율(0.0~1.0)로 직접 반환됨
                        crop_info = {
                            "left":   round(float(cl), 4),
                            "top":    round(float(ct), 4),
                            "right":  round(float(cr), 4),
                            "bottom": round(float(cb), 4),
                        }
                        is_cropped = any(v > 0.01 for v in crop_info.values())
                    except Exception:
                        pass

                    images.append({
                        "data":       img_blob,
                        "ext":        img_ext,
                        "w":          w,
                        "h":          h,
                        "pil":        pil_img,
                        "crop_info":  crop_info,
                        "is_cropped": is_cropped,
                    })
            except Exception:
                pass

            # 그룹 shape 안 텍스트 + 이미지 재귀 처리
            try:
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for child in _iter_group_shapes(shape):
                        # 텍스트박스
                        if child.has_text_frame:
                            for para in child.text_frame.paragraphs:
                                line = para.text or ""
                                for run in para.runs:
                                    try:
                                        if run.hyperlink and run.hyperlink.address:
                                            hyperlinks.append(run.hyperlink.address)
                                    except Exception:
                                        pass
                                if line.strip():
                                    text_items.append(("textbox", line.strip()))
                        # 이미지
                        if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            try:
                                img_blob = child.image.blob
                                img_ext  = child.image.ext
                                pil_img  = Image.open(io.BytesIO(img_blob))
                                w, h     = pil_img.size
                                images.append({
                                    "data": img_blob,
                                    "ext":  img_ext,
                                    "w":    w,
                                    "h":    h,
                                    "pil":  pil_img,
                                })
                            except Exception:
                                pass
            except Exception:
                pass

        # ── 슬라이드 노트 ────────────────────────────────────
        try:
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                for para in notes_tf.paragraphs:
                    nt = para.text.strip()
                    if nt:
                        text_items.append(("note", nt))
        except Exception:
            pass

        # 전체 텍스트 합산
        full_text = "\n".join(t for _, t in text_items)

        return SlideData(
            slide_number=slide_num,
            text=full_text,
            text_items=text_items,
            images=images,
            hyperlinks=list(set(hyperlinks)),
            is_hidden=is_hidden,
        )

    # ── 슬라이드 썸네일 (이미지 → PIL) ───────────────────────
    def slide_thumbnail(self, idx: int, max_w: int = 1000) -> Optional[Image.Image]:
        """
        슬라이드 전체를 하나의 이미지로 렌더링.
        LibreOffice 없이 python-pptx만으로는 불가하므로
        슬라이드 내 이미지들을 흰 캔버스에 합성하여 근사치 렌더링.
        → Claude Vision AI 입력용
        """
        try:
            slide_data = self.extract_slide(idx)
            if not slide_data.images:
                return None

            # 슬라이드 크기 (EMU → px, 96dpi 기준)
            prs = self._prs
            slide_w_px = int(prs.slide_width  / 914400 * 96)
            slide_h_px = int(prs.slide_height / 914400 * 96)

            # 스케일 조정
            scale = min(max_w / max(slide_w_px, 1), 1.5)
            canvas_w = int(slide_w_px * scale)
            canvas_h = int(slide_h_px * scale)

            canvas = Image.new("RGB", (canvas_w, canvas_h), (255, 255, 255))

            # 이미지들을 캔버스 크기에 맞게 배치 (단순 타일링)
            slide      = prs.slides[idx]
            img_shapes = []
            try:
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        img_shapes.append(shape)
                    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        for child in shape.shapes:
                            if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                img_shapes.append(child)
            except Exception:
                pass

            for shape in img_shapes:
                try:
                    left = int(shape.left   / 914400 * 96 * scale)
                    top  = int(shape.top    / 914400 * 96 * scale)
                    w    = int(shape.width  / 914400 * 96 * scale)
                    h    = int(shape.height / 914400 * 96 * scale)

                    pil = Image.open(io.BytesIO(shape.image.blob))
                    if pil.mode not in ("RGB", "RGBA"):
                        pil = pil.convert("RGB")
                    if w > 0 and h > 0:
                        pil = pil.resize((w, h), Image.LANCZOS)
                    canvas.paste(pil, (max(0, left), max(0, top)))
                except Exception:
                    pass

            return canvas
        except Exception as e:
            logger.warning(f"슬라이드 썸네일 실패 s{idx+1}: {e}")
            return None

    # ── 슬라이드 전체 이미지 렌더링 (Claude Vision 입력) ──────
    def render_for_vision(self, idx: int) -> Optional[Image.Image]:
        """Claude Vision AI 입력용 슬라이드 이미지"""
        return self.slide_thumbnail(idx, max_w=1200)

    # ── Layout / Master 이미지 직접 추출 ─────────────────────
    def extract_layout_master_images(self) -> list[dict]:
        """
        슬라이드 레이아웃(slideLayout)과 슬라이드 마스터(slideMaster)에서
        이미지를 직접 추출합니다.

        반환: [
          {
            "data":   bytes,        # 원본 바이너리
            "ext":    str,          # 확장자 ("jpeg", "png", …)
            "w":      int,          # 픽셀 너비
            "h":      int,          # 픽셀 높이
            "pil":    Image,        # PIL 이미지
            "source": str,          # "layout" | "master"
            "source_idx": int,      # 레이아웃/마스터 0-based 인덱스
            "shape_name": str,      # shape.name (디버그용)
          },
          …
        ]
        중복 이미지(동일 blob)는 제거합니다.
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        results: list[dict] = []
        seen_blobs: set[bytes] = set()

        def _extract_from_shapes(shapes, source: str, source_idx: int):
            for shape in shapes:
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        blob = shape.image.blob
                        if blob in seen_blobs:
                            continue
                        seen_blobs.add(blob)
                        pil = Image.open(io.BytesIO(blob))
                        results.append({
                            "data":       blob,
                            "ext":        shape.image.ext,
                            "w":          pil.size[0],
                            "h":          pil.size[1],
                            "pil":        pil,
                            "source":     source,
                            "source_idx": source_idx,
                            "shape_name": getattr(shape, "name", ""),
                        })
                    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        # 그룹 내부 재귀
                        _extract_from_shapes(
                            list(_iter_group_shapes(shape)), source, source_idx
                        )
                except Exception:
                    pass

        try:
            prs = self._prs
            # ── 슬라이드 마스터 ──
            for mi, master in enumerate(prs.slide_masters):
                try:
                    _extract_from_shapes(master.shapes, "master", mi)
                except Exception as e:
                    logger.debug(f"master[{mi}] 추출 오류: {e}")

            # ── 슬라이드 레이아웃 ──
            for mi, master in enumerate(prs.slide_masters):
                for li, layout in enumerate(master.slide_layouts):
                    try:
                        _extract_from_shapes(layout.shapes, "layout", li)
                    except Exception as e:
                        logger.debug(f"layout[{mi}/{li}] 추출 오류: {e}")

        except Exception as e:
            logger.warning(f"layout/master 이미지 추출 실패: {e}")

        logger.info(
            f"layout/master 이미지 추출: {len(results)}개 "
            f"(master={sum(1 for r in results if r['source']=='master')}, "
            f"layout={sum(1 for r in results if r['source']=='layout')})"
        )
        return results

    # ── Layout / Master 텍스트 추출 ──────────────────────────
    def extract_layout_master_texts(self) -> list[dict]:
        """
        슬라이드 마스터(slideMaster)와 슬라이드 레이아웃(slideLayout)에서
        텍스트를 추출합니다.

        반환: [
          {
            "text":       str,   # 추출된 텍스트
            "source":     str,   # "master" | "layout"
            "source_idx": int,   # 마스터/레이아웃 0-based 인덱스
            "layout_name": str,  # 레이아웃 이름 (source==layout일 때)
            "affected_slides": list[int],  # 이 마스터/레이아웃이 적용된 슬라이드 번호 목록
          },
          ...
        ]
        - 동일 텍스트가 여러 마스터/레이아웃에 중복될 경우 한 번만 포함.
        - 슬라이드 본문에 이미 존재하는 텍스트는 포함하지 않음 (중복 방지는
          pipeline 단에서 처리).
        """
        results: list[dict] = []
        seen_texts: set[str] = set()  # 마스터/레이아웃 내부 중복 제거

        def _iter_shapes_text(shapes):
            """shape 목록에서 텍스트 재귀 수집 (GROUP 포함)"""
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            for shape in shapes:
                try:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                yield t
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                t = cell.text.strip()
                                if t:
                                    yield t
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        yield from _iter_shapes_text(shape.shapes)
                except Exception:
                    pass

        try:
            prs = self._prs
            master_list = list(prs.slide_masters)

            # 마스터→적용 슬라이드 매핑 구축
            master_to_slides: dict[int, list[int]] = {}
            layout_to_slides: dict[tuple, list[int]] = {}  # (master_idx, layout_idx) → [slide_num]
            slides_list = list(prs.slides)
            for si, slide in enumerate(slides_list):
                slide_num = si + 1
                try:
                    lo = slide.slide_layout
                    for mi, m in enumerate(master_list):
                        try:
                            li = list(m.slide_layouts).index(lo)
                            master_to_slides.setdefault(mi, []).append(slide_num)
                            layout_to_slides.setdefault((mi, li), []).append(slide_num)
                            break
                        except ValueError:
                            continue
                except Exception:
                    pass

            # ── 슬라이드 마스터 텍스트 ──
            for mi, master in enumerate(master_list):
                affected = master_to_slides.get(mi, [])
                try:
                    for t in _iter_shapes_text(master.shapes):
                        if t in seen_texts:
                            continue
                        seen_texts.add(t)
                        results.append({
                            "text":          t,
                            "source":        "master",
                            "source_idx":    mi,
                            "layout_name":   "",
                            "affected_slides": affected,
                        })
                except Exception as e:
                    logger.debug(f"master[{mi}] 텍스트 추출 오류: {e}")

            # ── 슬라이드 레이아웃 텍스트 ──
            for mi, master in enumerate(master_list):
                for li, layout in enumerate(master.slide_layouts):
                    affected = layout_to_slides.get((mi, li), [])
                    try:
                        for t in _iter_shapes_text(layout.shapes):
                            if t in seen_texts:
                                continue
                            seen_texts.add(t)
                            results.append({
                                "text":          t,
                                "source":        "layout",
                                "source_idx":    li,
                                "layout_name":   getattr(layout, "name", ""),
                                "affected_slides": affected,
                            })
                    except Exception as e:
                        logger.debug(f"layout[{mi}/{li}] 텍스트 추출 오류: {e}")

        except Exception as e:
            logger.warning(f"layout/master 텍스트 추출 실패: {e}")

        logger.info(
            f"layout/master 텍스트 추출: {len(results)}개 "
            f"(master={sum(1 for r in results if r['source']=='master')}, "
            f"layout={sum(1 for r in results if r['source']=='layout')})"
        )
        return results
