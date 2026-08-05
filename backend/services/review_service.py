"""
제안서 검수 서비스
- 4개 파일(감리RFP, 대상사업RFP, 포털HTML, 제안서PPT) → Claude 단일 호출 분석 → JSON 리포트
- 멀티턴 Tool Use 루프 제거: 4개 문서 전체 텍스트를 한 번에 입력하여 1회 호출로 완료
- 토큰 절감: 멀티턴 히스토리 누적 없음 (~90% 절감)
"""
from __future__ import annotations
import io
import json
import base64
import re
from pathlib import Path
from typing import Optional

from core.config import get_logger, ANTHROPIC_API_KEY, DATA_DIR

logger = get_logger("review_service")

# 검수 전용 모델 (관리자에서 변경 가능)
_REVIEW_MODEL_FILE = DATA_DIR / "review_model.txt"
DEFAULT_REVIEW_MODEL = "claude-haiku-4-5"


def get_review_model() -> str:
    try:
        if _REVIEW_MODEL_FILE.exists():
            m = _REVIEW_MODEL_FILE.read_text("utf-8").strip()
            if m:
                # 잘못된 날짜 suffix 자동 보정
                if m in ("claude-sonnet-4-5-20250514", "claude-sonnet-4-5"):
                    m = "claude-haiku-4-5"
                    _REVIEW_MODEL_FILE.write_text(m, "utf-8")
                    logger.info(f"[review] 모델 자동 교체: {m}")
                return m
    except Exception:
        pass
    return DEFAULT_REVIEW_MODEL


def set_review_model(model: str):
    try:
        _REVIEW_MODEL_FILE.write_text(model.strip(), "utf-8")
    except Exception as e:
        logger.warning(f"review model 저장 실패: {e}")


def _extract_text_from_hwp(data: bytes) -> str:
    """HWP 5.x (OLE 컨테이너) → 텍스트 추출"""
    try:
        import olefile
        import zlib
        import struct

        ole = olefile.OleFileIO(io.BytesIO(data))

        # FileHeader에서 압축 여부 플래그 확인 (offset 36, bit0)
        try:
            hdr = ole.openstream("FileHeader").read()
            is_compressed = bool(struct.unpack_from("<I", hdr, 36)[0] & 1)
        except Exception:
            is_compressed = True

        texts: list[str] = []
        for i in range(1, 300):
            stream_name = f"BodyText/Section{i:04d}"
            if not ole.exists(stream_name):
                break
            raw = ole.openstream(stream_name).read()
            if is_compressed:
                try:
                    raw = zlib.decompress(raw, -15)
                except Exception:
                    pass

            # HWP 레코드 스트림 파싱 — PARA_TEXT (type=67) 레코드만 추출
            pos = 0
            while pos + 4 <= len(raw):
                hword = struct.unpack_from("<I", raw, pos)[0]
                rec_type = hword & 0x3FF
                rec_size = (hword >> 20) & 0xFFF
                if rec_size == 0xFFF:
                    if pos + 8 > len(raw):
                        break
                    rec_size = struct.unpack_from("<I", raw, pos + 4)[0]
                    pos += 8
                else:
                    pos += 4
                payload = raw[pos: pos + rec_size]
                pos += rec_size

                if rec_type == 67:  # PARA_TEXT
                    try:
                        chars: list[str] = []
                        for j in range(0, len(payload) - 1, 2):
                            code = struct.unpack_from("<H", payload, j)[0]
                            if code in (0x000D, 0x0000):
                                chars.append("\n")
                            elif code >= 0x0020:
                                chars.append(chr(code))
                        line = "".join(chars).strip()
                        if line:
                            texts.append(line)
                    except Exception:
                        pass
        ole.close()
        result = "\n".join(texts)
        logger.debug(f"[review] HWP 추출: {len(result):,}자")
        return result
    except Exception as e:
        logger.warning(f"[review] HWP 텍스트 추출 실패: {e}")
        return ""


def _extract_text_from_hwpx(data: bytes) -> str:
    """HWPX (ZIP+XML 컨테이너) → 텍스트 추출"""
    try:
        import zipfile
        from lxml import etree

        texts: list[str] = []
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            # Contents/section0.xml, section1.xml, ...
            section_files = sorted(
                n for n in zf.namelist()
                if n.startswith("Contents/section") and n.endswith(".xml")
            )
            for fname in section_files:
                xml_bytes = zf.read(fname)
                try:
                    root = etree.fromstring(xml_bytes)
                    for elem in root.iter():
                        tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
                        if tag == "t" and elem.text and elem.text.strip():
                            texts.append(elem.text.strip())
                except Exception:
                    pass
        result = "\n".join(texts)
        logger.debug(f"[review] HWPX 추출: {len(result):,}자")
        return result
    except Exception as e:
        logger.warning(f"[review] HWPX 텍스트 추출 실패: {e}")
        return ""


def _extract_text_from_pdf(data: bytes) -> str:
    """PDF → 텍스트 추출 (pdfplumber 우선, 실패 시 PyPDF2)"""
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            pages = []
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(f"[페이지 {i+1}]\n{text}")
            return "\n\n".join(pages)
    except Exception as e1:
        logger.debug(f"pdfplumber 실패: {e1}")
        try:
            import PyPDF2
            reader = PyPDF2.PdfReader(io.BytesIO(data))
            pages = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append(f"[페이지 {i+1}]\n{text}")
            return "\n\n".join(pages)
        except Exception as e2:
            logger.warning(f"PDF 텍스트 추출 실패: {e2}")
            return ""


def _extract_text_from_pptx(data: bytes) -> str:
    """PPTX → 텍스트 추출 (슬라이드별)

    ⚠️ python-pptx shape.text_frame.text 만으로는 두 가지 텍스트가 누락된다:
      1) GroupShape 내부 shape — slide.shapes 순회로는 접근 불가
      2) run이 여러 <a:t>로 쪼개진 경우 — paragraph 내 run들이 분리 저장됨

    → slide._element XML을 직접 이터레이트해 <a:p>(paragraph) 단위로
      모든 <a:t> run을 concat하면 두 문제 모두 해결된다.
    표(table)는 <a:tc>(cell) 단위로 별도 처리한다.
    """
    try:
        from pptx import Presentation

        # DrawingML 네임스페이스
        _A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        _TAG_P  = f'{{{_A}}}p'   # paragraph
        _TAG_T  = f'{{{_A}}}t'   # text run
        _TAG_TC = f'{{{_A}}}tc'  # table cell
        _TAG_TR = f'{{{_A}}}tr'  # table row

        prs = Presentation(io.BytesIO(data))
        slide_texts: list[str] = []

        for i, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            # 표 cell XML 경로 — 중복 수집 방지용
            table_cell_elems: set[int] = set()

            # ── 1. 표(table) 먼저 수집 — 행 단위 | 구분으로 ────────
            for tc_parent in slide._element.iter(_TAG_TR):
                row_parts: list[str] = []
                for tc in tc_parent.iter(_TAG_TC):
                    table_cell_elems.add(id(tc))
                    cell_text = "".join(
                        t.text for t in tc.iter(_TAG_T) if t.text
                    ).strip()
                    row_parts.append(cell_text if cell_text else "-")
                if any(v != "-" for v in row_parts):
                    texts.append(" | ".join(row_parts))

            # ── 2. 표 외 나머지 paragraph 수집 ──────────────────────
            for p in slide._element.iter(_TAG_P):
                # 표 cell 안에 있는 paragraph는 이미 처리했으므로 건너뜀
                parent = p.getparent()
                if parent is not None and id(parent) in table_cell_elems:
                    continue
                # 같은 paragraph 안의 모든 run을 concat
                line = "".join(t.text for t in p.iter(_TAG_T) if t.text).strip()
                if line:
                    texts.append(line)

            if texts:
                slide_texts.append(f"[슬라이드 {i}]\n" + "\n".join(texts))

        return "\n\n".join(slide_texts)
    except Exception as e:
        logger.warning(f"PPTX 텍스트 추출 실패: {e}")
        return ""


def _extract_text_from_html(data: bytes) -> str:
    """HTML → 텍스트 추출 (테이블은 마크다운으로 변환하여 컬럼 구조 보존)

    ⚠️ 핵심: get_text()만 쓰면 테이블 셀이 모두 줄바꿈으로 나열됨.
    예) 예비조사(MD)=2, 감리(MD)=12, 조치확인(MD)=3, 제안(MD)=19 가 모두
        "2 12 3 19" 처럼 뭉쳐서 Claude가 어느 컬럼인지 판별 불가.
    → 테이블을 마크다운 | col | col | 형태로 변환하면 열 위치가 명확해짐.
    """
    try:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(data, "html.parser")
        # script/style 제거
        for tag in soup(["script", "style"]):
            tag.decompose()

        # ── 테이블 → 마크다운 변환 ─────────────────────────────────
        for table in soup.find_all("table"):
            rows = table.find_all("tr")
            if not rows:
                table.decompose()
                continue

            md_lines: list[str] = []
            for i, tr in enumerate(rows):
                cells = tr.find_all(["th", "td"])
                if not cells:
                    continue
                values = [c.get_text(separator=" ", strip=True).replace("|", "｜")
                          for c in cells]
                md_lines.append("| " + " | ".join(values) + " |")
                # 헤더 다음 구분선
                if i == 0 and tr.find("th"):
                    md_lines.append("| " + " | ".join(["---"] * len(values)) + " |")

            # 테이블 원소를 마크다운 텍스트 노드로 교체
            from bs4 import NavigableString
            table.replace_with(NavigableString("\n" + "\n".join(md_lines) + "\n"))

        return soup.get_text(separator="\n", strip=True)

    except Exception:
        # fallback: 정규식 (구조 무너지지만 최소한 텍스트는 유지)
        text = data.decode("utf-8", errors="ignore")
        text = re.sub(r"<script[^>]*>.*?</script>", "", text, flags=re.DOTALL)
        text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL)
        text = re.sub(r"<[^>]+>", " ", text)
        return re.sub(r"\s+", "\n", text).strip()


def _parse_portal_html(data: bytes) -> str:
    """포털 HTML에서 tblSchedule / tblManList 를 직접 파싱하여
    Claude가 숫자를 절대 틀리지 않도록 구조화된 텍스트로 반환.

    tblSchedule 구조:
      - 단계별 2행(감리원 행 + 전문가 행), rowspan으로 단계명·날짜 공유
      - 컬럼 순서: [단계명] [날짜] [인력구분] [인원수] [예비조사MD] [감리MD] [조치확인MD] [제안MD] [투입인력]
      - HTML공수(MD) = 같은 단계 감리원 제안MD + 전문가 제안MD
      - 맨 아래 합계 행으로 검증

    tblManList 구조:
      - 인력 구분(단계감리팀/전문가팀 등), 성명, 정기/추가/검수지원/소계 MD, 상근 여부, 등급
    """
    try:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(data, "html.parser")
        lines: list[str] = []

        # ── 1. tblSchedule 파싱 ────────────────────────────────────
        tbl_sched = soup.find("table", {"id": "tblSchedule"})
        if tbl_sched:
            lines.append("=== 포털 감리 일정 (tblSchedule) ===")
            lines.append("※ HTML공수(MD) = 각 단계의 [감리원 제안MD + 전문가 제안MD]")
            lines.append("")

            rows = tbl_sched.find_all("tr")

            # rowspan 처리용: 현재 행에 carryover될 (단계명, 날짜) 추적
            # {col_index: (value, remaining_rowspan)}
            carry: dict[int, tuple[str, int]] = {}

            # 헤더 행(th만 있는 행)은 건너뜀
            data_rows = [r for r in rows if r.find("td")]

            # 합계 행 분리 (th로만 구성된 합계 행이 마지막에 존재)
            # "합계" 텍스트를 포함하는 마지막 행
            sum_row = None
            for r in rows:
                if "합계" in r.get_text():
                    sum_row = r

            stage_data: list[dict] = []  # [{단계명, 날짜, 인력구분, 인원수, 예비조사, 감리, 조치확인, 제안MD, 투입인력}]

            for tr in data_rows:
                if tr == sum_row:
                    continue

                tds = tr.find_all(["td", "th"])
                # carry 적용: 빠진 앞 컬럼(rowspan으로 이전 행이 차지) 채워넣기
                full_cells: list[str] = []
                cell_iter = iter(tds)

                # 최대 9컬럼 (단계명, 날짜, 인력구분, 인원, 예비조사, 감리, 조치확인, 제안MD, 투입인력)
                col_idx = 0
                td_list = list(tds)
                td_ptr = 0

                result_cols: list[str] = []
                for col_idx in range(9):
                    if col_idx in carry:
                        val, rem = carry[col_idx]
                        result_cols.append(val)
                        if rem - 1 > 0:
                            carry[col_idx] = (val, rem - 1)
                        else:
                            del carry[col_idx]
                    else:
                        if td_ptr >= len(td_list):
                            result_cols.append("")
                            continue
                        td = td_list[td_ptr]
                        td_ptr += 1
                        val = td.get_text(separator=" ", strip=True)
                        rs = int(td.get("rowspan", 1))
                        if rs > 1:
                            carry[col_idx] = (val, rs - 1)
                        result_cols.append(val)

                # 컬럼 매핑
                if len(result_cols) >= 8:
                    stage_data.append({
                        "단계명":   result_cols[0],
                        "날짜":     result_cols[1],
                        "인력구분": result_cols[2],
                        "인원수":   result_cols[3],
                        "예비조사": result_cols[4],
                        "감리MD":   result_cols[5],
                        "조치확인": result_cols[6],
                        "제안MD":   result_cols[7],
                        "투입인력": result_cols[8] if len(result_cols) > 8 else "",
                    })

            # 단계별로 그룹핑하여 출력
            stages_seen: list[str] = []
            stage_groups: dict[str, list[dict]] = {}
            for row in stage_data:
                sname = row["단계명"]
                if sname not in stage_groups:
                    stages_seen.append(sname)
                    stage_groups[sname] = []
                stage_groups[sname].append(row)

            total_html_md = 0
            for sname in stages_seen:
                group = stage_groups[sname]
                감리원_md = 0
                전문가_md = 0
                날짜 = ""
                for row in group:
                    날짜 = row["날짜"] or 날짜
                    try:
                        md_val = int(row["제안MD"])
                    except (ValueError, TypeError):
                        md_val = 0
                    if "감리원" in row["인력구분"]:
                        감리원_md = md_val
                    elif "전문가" in row["인력구분"]:
                        전문가_md = md_val
                단계_md = 감리원_md + 전문가_md
                total_html_md += 단계_md
                lines.append(f"[{sname}]")
                lines.append(f"  날짜: {날짜}")
                lines.append(f"  감리원 제안MD: {감리원_md}")
                lines.append(f"  전문가 제안MD: {전문가_md}")
                lines.append(f"  단계 합계MD:   {단계_md}  ← HTML공수(MD)")
                lines.append("")

            # 합계 행 파싱
            if sum_row:
                sum_cells = sum_row.find_all(["th", "td"])
                sum_vals = [c.get_text(strip=True) for c in sum_cells]
                # 합계 행: 합계 | 인원수 | 예비조사 | 감리 | 조치확인 | 제안MD | ...
                # 제안MD는 index 5 (0-base: 합계,인원수,예비조사,감리,조치확인,제안MD)
                portal_total = ""
                if len(sum_vals) >= 6:
                    portal_total = sum_vals[5]
                lines.append(f"[합계]")
                lines.append(f"  포털 합계 제안MD(합계행): {portal_total}")
                lines.append(f"  백엔드 직접 계산 합계MD:  {total_html_md}")
                lines.append(f"  {'✅ 일치' if str(total_html_md) == portal_total else '⚠️ 불일치 — 포털 합계행 값을 우선 사용'}")
                lines.append("")

        # ── 2. tblManList 파싱 ────────────────────────────────────
        tbl_man = soup.find("table", {"id": "tblManList"})
        if tbl_man:
            lines.append("=== 포털 제안 인력 (tblManList) ===")
            lines.append("구분 | 성명 | 정기MD | 추가MD | 검수지원MD | 소계MD | 상근여부 | 등급")
            lines.append("")

            man_rows = tbl_man.find_all("tr")
            # man_carry: {col_idx: (value, remaining_rowspan)}
            man_carry: dict[int, tuple[str, int]] = {}

            for tr in man_rows:
                tds = tr.find_all("td")
                if not tds:
                    continue

                # colspan/rowspan을 모두 고려해 13컬럼 배열로 펼치기
                # 컬럼 레이아웃(헤더 기준):
                # 0:구분  1:담당분야(colspan=2→1,2)  3:성명  4:정기  5:추가
                # 6:검수지원  7:소계  8:상근  9:등급  10:감리원증  11:연락처  12:교육시간
                result_cols: list[str] = []
                td_ptr = 0
                for col_idx in range(13):
                    if col_idx in man_carry:
                        val, rem = man_carry[col_idx]
                        result_cols.append(val)
                        if rem - 1 > 0:
                            man_carry[col_idx] = (val, rem - 1)
                        else:
                            del man_carry[col_idx]
                    else:
                        if td_ptr >= len(tds):
                            result_cols.append("")
                            continue
                        td = tds[td_ptr]
                        td_ptr += 1
                        val = td.get_text(separator=" ", strip=True)
                        rs = int(td.get("rowspan", 1))
                        cs = int(td.get("colspan", 1))
                        # rowspan: 다음 행들에 이 값 carry
                        if rs > 1:
                            man_carry[col_idx] = (val, rs - 1)
                        result_cols.append(val)
                        # colspan: 차지하는 나머지 컬럼도 같은 값으로 채움
                        for extra in range(1, cs):
                            next_col = col_idx + extra
                            result_cols.append(val)
                            if rs > 1:
                                man_carry[next_col] = (val, rs - 1)

                # 컬럼 인덱스: 0구분 1담당분야 2담당분야(colspan) 3성명 4정기 5추가 6검수지원 7소계 8상근 9등급
                if len(result_cols) >= 8:
                    구분    = result_cols[0]
                    성명    = result_cols[3] if len(result_cols) > 3 else ""
                    정기    = result_cols[4] if len(result_cols) > 4 else ""
                    추가    = result_cols[5] if len(result_cols) > 5 else ""
                    검수지원 = result_cols[6] if len(result_cols) > 6 else ""
                    소계    = result_cols[7] if len(result_cols) > 7 else ""
                    상근    = result_cols[8] if len(result_cols) > 8 else ""
                    등급    = result_cols[9] if len(result_cols) > 9 else ""
                    # 합계 행(성명 없고 숫자만 있는 행) 등 의미없는 행 건너뜀
                    if 성명 and not 성명.isdigit():
                        lines.append(f"  {구분} | {성명} | 정기={정기} 추가={추가} 검수지원={검수지원} 소계={소계} | {상근} | {등급}")

            lines.append("")

        result = "\n".join(lines)
        if not result.strip():
            # tblSchedule/tblManList 없으면 기존 방식 fallback
            return _extract_text_from_html(data)
        return result

    except Exception as e:
        logger.warning(f"포털 HTML 직접 파싱 실패, fallback: {e}")
        return _extract_text_from_html(data)




_SYSTEM_PROMPT = """\
당신은 정보시스템 감리사업 제안서 전문 검수 AI입니다.
사용자가 제공하는 4개 문서를 분석하여 정해진 JSON 스키마를 정확히 출력합니다.

검수 대상: 정보시스템 감리사업 정성제안서 (PPT)
검수 목적: 감리사업 RFP·포털 확정값과 제안서의 정합성 확인
검수 범위: 일정·공수·인력·잔존문구·오타 — 그 이상은 다루지 않는다

## 포털 HTML 공수(MD) 읽기 규칙 (반드시 준수)
- portal 문서의 각 단계에는 [감리원 제안MD], [전문가 제안MD], [단계 합계MD] 3줄이 있다
- **HTML공수(MD) = 반드시 "단계 합계MD" 값 사용** (감리원+전문가 합산값)
- 감리원 제안MD만 단독으로 HTML공수로 쓰면 절대 안 된다
- schedule JSON의 "HTML공수(MD)" 컬럼 = 각 단계의 "단계 합계MD" 값

## 절대 규칙
1. 응답은 반드시 유효한 JSON 객체 **하나만** 출력한다. 코드블록(```), 설명 문구, 마크다운 일절 금지.
2. JSON 문자열 값 안에 실제 줄바꿈 문자(0x0A/0x0D)를 절대 사용하지 않는다. 줄바꿈이 필요하면 반드시 \\n 이스케이프 시퀀스를 사용한다.
3. 모든 문자열 내 이중인용부호(")는 반드시 \\" 로 이스케이프한다.
4. 배열·객체 값이 없을 때는 null 대신 빈 배열 [] 또는 빈 문자열 ""을 사용한다.

## 출력 JSON 스키마 (키와 타입을 정확히 지킬 것 — 스키마에 없는 키는 출력하지 않는다)
{
  "id": "string — 영문소문자-숫자-하이픈 슬러그",
  "name": "string — 감리사업명(감리사업 RFP 기준)",
  "org": "string — 발주기관명(감리사업 RFP 기준)",
  "date": "string — 검수일",
  "counts": {"crit": 0, "major": 0, "minor": 0, "check": 0},
  "verdict": "string — 슬라이드별 수정 권고사항 취합. 형식: 슬라이드 번호 오름차순으로 한 줄씩, 각 줄은 '<b>슬라이드 N</b> 수정 내용' 형식. 줄바꿈은 \\n. 오류/확인 항목이 없는 슬라이드는 생략. 마지막 줄은 전체 요약 한 줄로 마무리.",
  "baseline": [["항목명","기준값","출처"], ...],
  "critical": [{"title":"string","slide":"string","fix":"string","body":"string"}, ...],
  "major":    [{"title":"string","slide":"string","fix":"string","body":"string"}, ...],
  "minor":    [{"title":"string","slide":"string","fix":"string","body":"string"}, ...],
  "checkNeeded": [{"title":"string","slide":"string","fix":"string","body":"string"}, ...],
  "schedule": [["단계명","HTML기준일정","PPT일정","HTML공수(MD)","PPT공수(MD)","ok|major|check|crit"], ...],
  // ⚠️ schedule 배열에는 단계 행만 넣는다. "합계" 행은 절대 포함하지 않는다. 합계는 프론트엔드가 자동 계산한다.
  "scheduleNote": "string",
  "personnel": [["구분","HTML기준인력","PPT표기인력","ok|major|check|crit","슬라이드번호","불일치내용(ok이면 빈 문자열)"], ...],
  "irrelevant": {"summary":"string","items":[{"text":"string","slide":"string","sourceGuess":"string","severity":"high|low"}]},
  "typoChecklist": [{"no":1,"slide":"string","type":"string","priority":"높음|중간|낮음","original":"string","fix":"string","note":"string"}, ...],
  "typoNote": "string",
  "priority": {"crit":["string"],"major":["string"],"check":["string"]}
}

## 검수 규율 (반드시 지킬 것 — 위반 시 그 항목은 아예 출력하지 않는다)

이 서비스는 정보시스템 감리사업 제안서를 10건 이상 실전 검수한 경험을 바탕으로
구축되었다. 개수 제한은 없으나 명확한 근거가 있는 항목만 기재한다.
근거 없이 항목을 만들어내지 않는다.
typoChecklist(오탈자)는 발견한 것을 개수와 무관하게 전부 기록한다.

검수 범위를 벗어난 항목은 출력하지 않는다:
- 범위 내: 일정·공수 대조, 인력명 대조, 잔존문구 검출, 오타·표기 오류
- 범위 밖(출력 금지): 예산·기성금·비용 계산, 과업 범위 커버리지, 자격 요건 세부 검토

### 규칙 A — 숫자/비율 관련 지적은 계산 없이는 절대 쓰지 않는다
숫자 불일치를 지적하려면:
1) 기준 문서에서 해당 수치를 직접 인용한다 (예: "포털 확정 총 공수 = 307 MD").
2) PPT에서 해당 수치를 직접 인용한다 (예: "슬라이드 122 표 합계 = 244 MD").
3) 두 값의 차이를 명시한다 (예: "307 - 244 = 63 MD 차이").
위 3단계 계산 없이는 숫자 관련 항목을 절대 critical/major/minor에 넣지 않는다.

### 규칙 B — 문장이 부자연스럽거나 잘려 보이면 그 자체를 오류로 쓰지 않는다
PPT 텍스트 추출 과정에서 발생한 깨짐·잘림·표 열 밀림은 실제 슬라이드 오류가 아닐 수 있다.
"문장이 끊겨 있다", "글자가 이상하다", "표가 깨졌다"는 이유만으로 오류로 보고하지 않는다.
실제 내용 불일치가 있는 경우에만 지적한다.
⚠️ 단, 명백한 오타(글자 오입력)·표기 혼용·외래어 오표기는 텍스트 추출 문제와 무관한 실제 오류다 — 규칙 B의 면제를 받지 않으며 typoChecklist에 반드시 기록한다.

### 규칙 C — "~해 보인다", "~일 가능성", "~인 것 같다" 같은 추측성 표현이 들어가면 그 항목을 삭제한다
body나 fix에 추측성 표현이 포함된 항목은 출력하지 않는다.
확실한 근거가 있는 경우에만 기재하며, 표현도 단정적으로 쓴다.

### 규칙 D — 근거 없는 항목은 넣지 않는다
카테고리별 개수 상한은 없다. 단, 명확한 근거가 있는 항목만 기재한다.
근거가 불분명한 항목은 개수와 무관하게 넣지 않는다.

### 규칙 E — 일정 관련 오탐을 미리 걸러낸다
다음 항목들은 오탐(false positive)이 잦으므로 특히 주의한다:
- 일정 표기 차이: PPT의 "YYYY.MM" vs HTML의 "YYYY-MM-DD"는 형식 차이일 뿐, 날짜 불일치로 보지 않는다.
- 공수(MD) 합산: 구성원별 공수를 직접 더해서 총합과 비교하지 않으면 지적하지 않는다.
- 인력 직함/등급 차이: RFP에 정확한 직함이 명시된 경우에만 지적하고, 그렇지 않으면 무시한다.

### 규칙 F — 대상사업 RFP는 irrelevant 전용이다
[대상사업 RFP]는 감리 대상인 SI사업의 발주 문서다.
이 문서는 **오직 irrelevant(잔존문구 검출)에서만** 참조한다:
- PPT 내용 중 대상사업과 무관한 **타 사업**의 업무범위·산출물·조직도가 복붙된 경우를 검출
- 단, 제안사의 수행실적·회사소개에 타사업명이 등장하는 것은 정상 — 잔존문구로 보지 않는다
- 발견이 없어도 irrelevant.summary는 반드시 작성한다

### 규칙 G — 예산·비용·기성금·과업범위·자격요건은 검수하지 않는다
이 항목들은 검수 범위 밖이다. 관련 내용이 보여도 critical/major/minor/checkNeeded에 넣지 않는다:
- 예산 총액, 기성금 비율(30/20/30/20 등), 부가세 포함 여부
- 과업 범위 커버리지(RFP 항목이 PPT에 반영되었는지)
- 감리원 자격 요건 세부 충족 여부

### verdict 작성 규칙 — 슬라이드별 수정 권고사항 취합
verdict는 검수에서 발견한 모든 수정 사항을 슬라이드 번호 오름차순으로 한 줄씩 나열한다.
- 형식: `<b>슬라이드 N</b> [심각도] 수정 내용`
  - 심각도 표기: 치명=❌, 중대=⚠️, 경미=🔹, 확인필요=❓
  - 예: `<b>슬라이드 12</b> ⚠️ 공수 합계 307MD → PPT 표기 244MD로 불일치, 수정 필요`
  - 예: `<b>슬라이드 45</b> 🔹 "정보시스템" → "정보 시스템" 오타`
- critical/major/minor/checkNeeded/typoChecklist의 모든 항목을 슬라이드 번호 기준으로 통합하여 나열
- 슬라이드 번호 없는 항목(slide가 빈 문자열)은 맨 앞에 `<b>전체</b> [내용]` 형식으로 기재
- 마지막 줄: `\n총 N건 (치명 a건 / 중대 b건 / 경미 c건 / 확인 d건)` 형식으로 요약
- 오류가 하나도 없으면: `이상 없음. 모든 항목이 기준과 일치합니다.`
- 줄바꿈은 반드시 \\n (실제 줄바꿈 절대 금지)

### 최종 자기점검 — 출력 직전 반드시 수행
JSON을 완성한 뒤, 출력하기 전에 다음을 점검한다:
1. 각 항목에 명확한 근거(출처 인용·수치 계산·슬라이드 번호)가 있는지 확인한다. 근거가 없는 항목은 삭제한다. 개수 제한은 없다.
2. 각 항목의 body에 규칙 A(계산 근거), 규칙 C(단정적 표현)를 위반한 것이 있으면 삭제한다.
3. counts 값은 최종 배열의 실제 길이와 일치해야 한다. 배열을 삭제했다면 counts도 갱신한다.
4. verdict는 위 "verdict 작성 규칙"에 따라 슬라이드별 수정 권고사항이 모두 포함되어 있는지 확인한다."""



# ── 오탈자 전용 독립 검출 함수 ────────────────────────────────────────────────
_TYPO_CHUNK_SIZE = 20  # 한 번에 처리할 슬라이드 수

def _detect_typos_standalone(ppt_text: str, api_key: str, model: str) -> list[dict]:
    """PPT 텍스트를 청크로 나눠 각각 독립 Claude 호출로 오탈자를 검출한다.
    메인 Tool Use 루프와 완전히 분리되어 히스토리 누적이 없다.
    반환: typoChecklist 형식의 dict 리스트
    """
    import re as _re
    import anthropic as _anthropic

    if not ppt_text.strip():
        return []

    # [슬라이드 N] 구분자 기준으로 분할
    parts = _re.split(r'(\[슬라이드 \d+\])', ppt_text)
    slides: list[tuple[int, str]] = []  # (슬라이드번호, 텍스트)
    i = 0
    while i < len(parts):
        m = _re.match(r'\[슬라이드 (\d+)\]', parts[i])
        if m:
            num = int(m.group(1))
            content = parts[i + 1] if i + 1 < len(parts) else ""
            slides.append((num, parts[i] + content))
            i += 2
        else:
            i += 1

    if not slides:
        return []

    total = slides[-1][0]
    cl = _anthropic.Anthropic(api_key=api_key)
    all_typos: list[dict] = []
    no_counter = [0]  # 전체 번호 누적용

    # 청크 단위로 독립 호출
    for chunk_start in range(0, len(slides), _TYPO_CHUNK_SIZE):
        chunk = slides[chunk_start: chunk_start + _TYPO_CHUNK_SIZE]
        chunk_text = "\n".join(c for _, c in chunk)
        slide_nums = [n for n, _ in chunk]
        s_from, s_to = slide_nums[0], slide_nums[-1]

        prompt = f"""다음은 정성제안서 PPT의 [슬라이드 {s_from}] ~ [슬라이드 {s_to}] 텍스트입니다 (전체 {total}슬라이드 중).
오탈자·표기 오류를 찾아 JSON 배열로만 반환하라. 설명·마크다운 일절 금지.

검출 대상:
- 명백한 오타: 글자 오입력 (예: 젋고→젊고, 재방방지→재발방지, 미관정보→민간정보)
- 단어 잘림: 문장 중간 단어 누락 (예: "효과성 극대"→"효과성 극대화")
- 표기 혼용: 같은 단어 다른 표기 (예: 어플리케이션↔애플리케이션, 워크샵↔워크숍)
- 외래어 오표기: 국립국어원 기준 위반 (예: 콘트롤러→컨트롤러)
- 맞춤법: 완전률→완전율, 률/율 구분 등
- 붙여쓰기: 고유명사·시스템명 공백 누락

반환 형식 (오탈자 없으면 빈 배열 [] 반환):
[{{"slide":"슬라이드번호","type":"오타|단어잘림|표기혼용|외래어|맞춤법|붙여쓰기","priority":"높음|중간|낮음","original":"원문","fix":"수정안","note":"설명"}}]

텍스트:
{chunk_text}"""

        try:
            resp = cl.messages.create(
                model=model,
                max_tokens=2000,
                messages=[{"role": "user", "content": prompt}],
            )
            raw = resp.content[0].text.strip() if resp.content else "[]"
            # JSON 배열 추출
            m2 = _re.search(r'\[.*\]', raw, _re.DOTALL)
            if m2:
                import json as _json
                chunk_typos = _json.loads(m2.group())
                for t in chunk_typos:
                    no_counter[0] += 1
                    t["no"] = no_counter[0]
                    all_typos.append(t)
            logger.info(f"[typo] 슬라이드 {s_from}~{s_to}: {len(chunk_typos) if m2 else 0}건 발견")
        except Exception as e:
            logger.warning(f"[typo] 슬라이드 {s_from}~{s_to} 오탈자 검출 실패: {e}")

    return all_typos

# ── Tool 정의: submit_report 단독 (단일 호출 방식) ──────────────────────────
TOOLS = [
    {
        "name": "submit_report",
        "description": (
            "검수를 완료하고 최종 JSON 보고서를 제출한다. "
            "4개 문서를 모두 분석한 뒤 이 도구를 한 번만 호출한다."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "report": {
                    "type": "object",
                    "description": "시스템 프롬프트 스키마를 따르는 최종 검수 JSON",
                    "properties": {
                        "id":    {"type": "string"},
                        "name":  {"type": "string"},
                        "org":   {"type": "string"},
                        "date":  {"type": "string"},
                        "counts": {
                            "type": "object",
                            "properties": {
                                "crit":  {"type": "integer"},
                                "major": {"type": "integer"},
                                "minor": {"type": "integer"},
                                "check": {"type": "integer"}
                            },
                            "required": ["crit", "major", "minor", "check"]
                        },
                        "verdict":       {"type": "string"},
                        "baseline":      {"type": "array"},
                        "critical":      {"type": "array"},
                        "major":         {"type": "array"},
                        "minor":         {"type": "array"},
                        "checkNeeded":   {"type": "array"},
                        "schedule":      {"type": "array"},
                        "scheduleNote":  {"type": "string"},
                        "personnel":     {"type": "array"},
                        "irrelevant":    {"type": "object"},
                        "typoChecklist": {"type": "array"},
                        "typoNote":      {"type": "string"},
                        "priority":      {"type": "object"}
                    },
                    "required": ["id", "name", "org", "date", "counts", "verdict",
                                 "baseline", "critical", "major", "minor", "checkNeeded",
                                 "schedule", "scheduleNote", "personnel", "irrelevant",
                                 "typoChecklist", "typoNote", "priority"]
                }
            },
            "required": ["report"]
        }
    }
]


def run_review(
    audit_rfp_data: bytes,   audit_rfp_name: str,
    target_rfp_data: bytes,  target_rfp_name: str,
    portal_html_data: bytes, portal_html_name: str,
    proposal_ppt_data: bytes, proposal_ppt_name: str,
    api_key: str = "",
    job_id: str = "",
) -> dict:
    """
    4개 파일을 Claude 단일 호출로 분석하여 검수 JSON 반환.
    4개 문서 전체 텍스트를 한 번에 입력하고 submit_report 도구로 결과를 수신한다.
    멀티턴 히스토리 누적 없음 → 토큰 ~90% 절감.
    """
    key = api_key or ANTHROPIC_API_KEY
    if not key:
        raise ValueError("Claude API 키가 설정되지 않았습니다.")
    model = get_review_model()

    # ── 텍스트 추출 ──────────────────────────────────────────────
    logger.info(f"[review] 텍스트 추출 시작: {audit_rfp_name}, {target_rfp_name}, {portal_html_name}, {proposal_ppt_name}")

    def extract(data: bytes, name: str) -> str:
        ext = Path(name).suffix.lower()
        if ext == ".pdf":
            return _extract_text_from_pdf(data)
        elif ext in (".pptx", ".ppt"):
            return _extract_text_from_pptx(data)
        elif ext in (".html", ".htm"):
            return _parse_portal_html(data)
        elif ext == ".hwpx":
            return _extract_text_from_hwpx(data)
        elif ext == ".hwp":
            return _extract_text_from_hwp(data)
        else:
            return data.decode("utf-8", errors="ignore")

    audit_rfp_text    = extract(audit_rfp_data,    audit_rfp_name)
    target_rfp_text   = extract(target_rfp_data,   target_rfp_name)
    portal_html_text  = extract(portal_html_data,  portal_html_name)
    proposal_ppt_text = extract(proposal_ppt_data, proposal_ppt_name)

    logger.info(
        f"[review] 추출 완료: 감리RFP={len(audit_rfp_text):,}자 "
        f"대상RFP={len(target_rfp_text):,}자 "
        f"포털={len(portal_html_text):,}자 "
        f"PPT={len(proposal_ppt_text):,}자"
    )

    # ── 날짜 ─────────────────────────────────────────────────────
    from core.config import now_kst
    today = now_kst().strftime("%Y.%m.%d")

    # ── 오탈자 사전 검출 (메인 호출과 독립) ──────────────────────
    logger.info("[review] 오탈자 사전 검출 시작")
    typo_results = _detect_typos_standalone(proposal_ppt_text, key, model)
    logger.info(f"[review] 오탈자 사전 검출 완료: {len(typo_results)}건")

    import json as _json
    typo_json_str = _json.dumps(typo_results, ensure_ascii=False, indent=2)

    # ── 단일 호출용 사용자 메시지 구성 ───────────────────────────
    # 4개 문서 전체 텍스트를 한 번에 포함하여 1회 Claude 호출로 완료
    user_content = f"""오늘 날짜: {today}

아래 4개 문서를 분석하여 정성제안서 PPT를 검수하고, 즉시 submit_report 도구로 최종 JSON을 제출하라.
도구 호출은 submit_report 하나만 허용된다. 추가 검색 없이 아래 제공된 문서만으로 검수한다.

## 포털 HTML 공수(MD) / 일정 읽기 방법
- portal 문서의 각 단계: [단계명] / 날짜 / 감리원 제안MD / 전문가 제안MD / **단계 합계MD** 순으로 명시
- **HTML공수(MD) = "단계 합계MD" 값만 사용** (감리원 단독값 사용 금지)
- 맨 아래 [합계] 섹션 전체 합계MD로 검증

## PPT 공수(MD) 읽기 규칙
- **PPT공수(MD) = 해당 단계 투입 모든 인력(감리원+전문가+테스트팀 등) 공수 합계**
- 감리원(단계 감리팀) 공수만 쓰면 틀린다 — 전문가팀·테스트팀 반드시 포함
- "투입공수 합계", "총 XXX MD" 합계 셀이 있으면 우선 활용

## 대상사업 RFP 사용 범위
- [대상사업 RFP]는 **irrelevant(잔존문구 검출) 전용** — 다른 검수 항목에 사용하지 않는다

## 오탈자(typoChecklist) — 사전 검출 결과
오탈자 검출은 이미 완료되었다. 아래 결과를 typoChecklist에 그대로 반영하고 추가 발견분을 더한다.
{typo_json_str}

## 검수 항목 (9개 — 모두 포함)
baseline / critical / major / minor / checkNeeded /
schedule+scheduleNote / personnel / irrelevant / typoChecklist+typoNote

---

## [감리사업 RFP] — 사업명·발주기관·일정·인력 기준
{audit_rfp_text}

---

## [대상사업 RFP] — irrelevant 검출 전용
{target_rfp_text}

---

## [포털 제안작업표 HTML] — 일정·공수·인력 확정값
{portal_html_text}

---

## [정성제안서 PPT] — 검수 대상
{proposal_ppt_text}

---

위 4개 문서를 바탕으로 9개 검수 항목을 모두 작성하여 즉시 submit_report로 제출하라."""

    # ── Claude 단일 호출 ──────────────────────────────────────────
    import anthropic

    cl = anthropic.Anthropic(api_key=key)
    final_report: dict | None = None
    stop_reason = "unknown"

    logger.info(f"[review] 단일 호출 시작: model={model}, 입력={len(user_content):,}자")

    try:
        # 10분 초과 가능 요청 → 스트리밍 필수 (Anthropic SDK 정책)
        with cl.messages.stream(
            model=model,
            max_tokens=50000,
            system=_SYSTEM_PROMPT,
            tools=TOOLS,
            tool_choice={"type": "any"},   # submit_report 반드시 호출하도록 강제
            messages=[{"role": "user", "content": user_content}],
        ) as stream:
            response = stream.get_final_message()

        stop_reason = response.stop_reason
        in_tok  = response.usage.input_tokens  if response.usage else 0
        out_tok = response.usage.output_tokens if response.usage else 0
        logger.info(
            f"[review] 단일 호출 완료: stop={stop_reason}, "
            f"블록수={len(response.content)}, "
            f"입력={in_tok:,}tok 출력={out_tok:,}tok"
        )

        # submit_report 블록 추출
        for blk in response.content:
            name = getattr(blk, "name", None) or (blk.get("name") if isinstance(blk, dict) else None)
            if name == "submit_report":
                inp = getattr(blk, "input", None) or (blk.get("input") if isinstance(blk, dict) else {})
                final_report = inp.get("report", {})
                stop_reason = "submit_report"
                logger.info("[review] submit_report 수신 완료")
                break

        # submit_report 없이 end_turn으로 끝난 경우 → 텍스트에서 JSON 추출 시도
        if final_report is None:
            raw = ""
            for blk in response.content:
                t = getattr(blk, "type", None) or (blk.get("type") if isinstance(blk, dict) else None)
                if t == "text":
                    raw = getattr(blk, "text", "") or (blk.get("text", "") if isinstance(blk, dict) else "")
                    break
            logger.warning(f"[review] submit_report 없음 — 텍스트 JSON 추출 시도: {len(raw):,}자")

            def _sanitize_json_strings(s: str) -> str:
                res: list[str] = []
                in_string = False
                escaped = False
                for ch in s:
                    if escaped:
                        res.append(ch); escaped = False; continue
                    if ch == "\\":
                        escaped = True; res.append(ch); continue
                    if ch == '"':
                        in_string = not in_string; res.append(ch); continue
                    if in_string:
                        if ch == "\n": res.append("\\n")
                        elif ch == "\r": res.append("\\r")
                        elif ch == "\t": res.append("\\t")
                        else: res.append(ch)
                    else:
                        res.append(ch)
                return "".join(res)

            def _try_parse(s: str) -> dict | None:
                try:
                    return json.loads(s)
                except json.JSONDecodeError:
                    return None

            json_str = raw.strip()
            m2 = re.search(r"```(?:json)?\s*([\s\S]+?)\s*```", json_str)
            if m2:
                json_str = m2.group(1) if m2.group(1).startswith("{") else "{" + m2.group(1)
            result_candidate = _try_parse(json_str)
            if result_candidate is None:
                sanitized = _sanitize_json_strings(json_str)
                result_candidate = _try_parse(sanitized)
            if result_candidate is None:
                last_brace = json_str.rfind("}")
                if last_brace > 0:
                    result_candidate = _try_parse(json_str[:last_brace + 1]) or \
                                       _try_parse(_sanitize_json_strings(json_str[:last_brace + 1]))
            if result_candidate is not None:
                final_report = result_candidate
                logger.info("[review] 텍스트 JSON 추출 성공")
            else:
                logger.warning(f"[review] JSON 파싱 최종 실패. 처음 1000자: {raw[:1000]}")

    except Exception as e:
        logger.error(f"[review] 단일 호출 실패: {e}")
        stop_reason = f"error: {e}"

    # ── 결과 확정 ─────────────────────────────────────────────────
    if final_report is not None:
        result = final_report
        logger.info("[review] 결과 확정 완료")
    else:
        result = {
            "id": "parse-error",
            "name": "파싱 오류",
            "org": "",
            "date": f"{today} 검수",
            "counts": {"crit": 0, "major": 0, "minor": 0, "check": 0},
            "verdict": (
                "단일 호출 검수에서 submit_report가 반환되지 않았습니다. "
                f"종료 사유: {stop_reason}. "
                "파일을 다시 업로드하거나 관리자에게 문의하세요."
            ),
            "baseline": [],
            "critical": [], "major": [], "minor": [], "checkNeeded": [],
            "schedule": [], "scheduleNote": "",
            "personnel": [],
            "irrelevant": {"summary": "파싱 오류로 분석 불가", "items": []},
            "typoChecklist": [], "typoNote": "",
            "priority": {"crit": [], "major": [], "check": []},
            "_stop_reason": stop_reason,
        }

    # ── 필드 기본값 보정 ──────────────────────────────────────────
    for _removed in ("overview", "scopeCoverage", "qualificationCheck", "costCheck", "deliverableCheck"):
        result.pop(_removed, None)
    irr = result.get("irrelevant", "")
    if isinstance(irr, str):
        result["irrelevant"] = {
            "summary": irr if irr else "검출 결과 없음",
            "items": [],
        }
    else:
        result["irrelevant"].setdefault("summary", "")
        result["irrelevant"].setdefault("items", [])

    # counts 자동 계산 (Claude가 빠뜨린 경우 보정)
    result.setdefault("counts", {})
    result["counts"]["crit"]  = len(result.get("critical", []))
    result["counts"]["major"] = len(result.get("major", []))
    result["counts"]["minor"] = len(result.get("minor", []))
    result["counts"]["check"] = len(result.get("checkNeeded", []))

    return result
